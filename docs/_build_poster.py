"""Build the A0 NeuroMonster poster from the methodology-forward layout.

Produces docs/NoSegs_Poster_A0.pptx — fully editable native PowerPoint
(text boxes, table, picture objects). Figures inserted as pictures; all
text/guides/legends are editable. Run with the repo venv python.

Layout principles (2026-06 polish pass):
  * every major section opens with a short lead paragraph, then its figures;
  * figures are numbered Fig 1..13 in reading order (down col 1, then 2, 3);
  * key text upsized for A0 viewing distance;
  * vertical flow is estimated per block so columns do not overflow — each
    column prints its final y vs the panel bottom at build time.
"""
import math
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

REPO = Path(__file__).resolve().parent.parent
FIGS = REPO / "results" / "figures"
ASSETS = REPO / "docs" / "_poster_assets"
OUT = REPO / "docs" / "NoSegs_Poster_A0.pptx"

# ---- palette ----
PURPLE   = RGBColor(0x50, 0x07, 0x78)   # KCL purple, header/footer/bottom band
PANEL    = RGBColor(0xF3, 0xEE, 0xF7)   # light column panel
HEADING  = RGBColor(0x50, 0x07, 0x78)
SUBHEAD  = RGBColor(0x6A, 0x1B, 0x9A)
BODY     = RGBColor(0x2C, 0x2C, 0x2C)
META     = RGBColor(0xC9, 0xA8, 0xD4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LEGEND   = RGBColor(0x3A, 0x3A, 0x3A)
RULE     = RGBColor(0xC9, 0xA8, 0xD4)
F = "Arial"

# ---- type scale (upsized for A0) ----
SZ_BODY    = 18.5
SZ_LEAD    = 18.5
SZ_LEGEND  = 15.5
SZ_HEAD    = 32
SZ_SUBHEAD = 23
# middle column runs tighter on space, so its prose/figs are a touch smaller
SZ_BODY_M   = 16.5
SZ_LEGEND_M = 15.0

prs = Presentation()
prs.slide_width  = Cm(84.1)
prs.slide_height = Cm(118.9)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes


# ============================================================== helpers
def _no_line(sp):
    sp.line.fill.background()


def rect(x, y, w, h, fill, rounded=False, line=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = shapes.add_shape(shp, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if rounded:
        try:
            sp.adjustments[0] = 0.03
        except Exception:
            pass
    return sp


def textbox(x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
            wrap=True):
    """paras: list of dicts with keys:
       text, size, bold, italic, color, bullet(bool), align, space_after,
       runs (optional list of (text, dict-overrides)), line_spacing
    """
    tb = shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05); tf.margin_bottom = Cm(0.05)
    for i, pdef in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pdef.get("align", align)
        if pdef.get("space_after") is not None:
            p.space_after = Pt(pdef["space_after"])
        if pdef.get("space_before") is not None:
            p.space_before = Pt(pdef["space_before"])
        if pdef.get("line_spacing"):
            p.line_spacing = pdef["line_spacing"]
        runs = pdef.get("runs")
        if runs is None:
            runs = [(pdef.get("text", ""), {})]
        bullet = pdef.get("bullet", False)
        for txt, ov in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = ov.get("font", F)
            r.font.size = Pt(ov.get("size", pdef.get("size", 16)))
            r.font.bold = ov.get("bold", pdef.get("bold", False))
            r.font.italic = ov.get("italic", pdef.get("italic", False))
            r.font.color.rgb = ov.get("color", pdef.get("color", BODY))
        if bullet:
            _set_bullet(p)
    return tb


def _set_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set("indent", "-228600")
    pPr.set("marL", "228600")
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
    pPr.append(buFont)
    pPr.append(buChar)


# ---- vertical-flow estimation (so columns do not overflow) ----
def _para_text(pdef):
    if pdef.get("runs") is not None:
        return "".join(t for t, _ in pdef["runs"])
    return pdef.get("text", "")


def _para_h(text, size, w, line_spacing=1.04):
    # 0.0195 is deliberately conservative (Arial proportional + word-wrap slack)
    # so reserved height >= rendered height and blocks never overlap.
    cpl = max(1.0, w / (size * 0.0195))          # ~chars per line
    lines = max(1, math.ceil(len(text) / cpl))
    return lines * size * 0.0353 * line_spacing   # cm


def est_block(paras, w, default_size):
    h = 0.0
    for p in paras:
        size = p.get("size", default_size)
        h += _para_h(_para_text(p), size, w, p.get("line_spacing", 1.04))
        h += (p.get("space_after") or 0) / 28.35  # pt -> cm
    return h


def flow_text(x, w, y, paras, default_size, gap=0.35, pad=0.18):
    """Render a prose block and advance y by its estimated height."""
    h = est_block(paras, w, default_size) + pad
    textbox(x, y, w, h + 0.4, paras)
    return y + h + gap


def figure(path, x_col, colw, y, target_w=None, center=True, gap=0.25):
    """Place picture; width target_w (cm) or colw. Returns bottom y (cm)."""
    im = Image.open(path)
    ar = im.size[0] / im.size[1]
    w = target_w if target_w else colw
    h = w / ar
    x = x_col + (colw - w) / 2 if center else x_col
    shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    return y + h + gap


def _shape_text(sp, text, size=12.5, bold=True, color=WHITE, italic=False,
                align=PP_ALIGN.CENTER):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.15); tf.margin_right = Cm(0.15)
    tf.margin_top = Cm(0.04); tf.margin_bottom = Cm(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = F
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def flowchart(x, y, w):
    """Native, fully editable workflow flowchart with orange option chips and a
    quick-start panel. Returns bottom y (cm)."""
    TEAL     = RGBColor(0x1A, 0x7A, 0x7A)   # statistical-core highlight
    GREEN_IN = RGBColor(0x1A, 0x9D, 0x50)   # input stage
    BLUE_OUT = RGBColor(0x21, 0x66, 0xAC)   # output stage
    OBG = RGBColor(0xFF, 0xF4, 0xD6)        # orange chip fill
    OED = RGBColor(0xE0, 0xA8, 0x00)        # orange chip edge
    OTX = RGBColor(0x6B, 0x52, 0x00)        # orange chip text
    QBG = RGBColor(0xED, 0xE3, 0xF3)        # quick-start panel
    GRN = RGBColor(0x1A, 0x9D, 0x50)        # tick green

    bw, bh, gap = 13.4, 1.35, 0.55          # main box w/h + vertical gap (arrow)
    chip_x = x + bw + 0.9
    chip_w = w - bw - 0.9
    chip_h = 1.15
    cxc = x + bw / 2                         # arrow centre

    stages = [
        ("INPUT — raw transcripts  (x, y, gene, FOV)", GREEN_IN, None),
        ("1   PCA align + per-FOV GMM strips", PURPLE, "3-component GMM"),
        ("2   DBSCAN noise removal", PURPLE, "adaptive ε = p97 1-NN, [20, 30] px"),
        ("3   Observation window", PURPLE, "rectangle / convex / concave hull"),
        ("4   Bivariate Ripley K → L(r)", TEAL, "multi-scale r-grid (cell + tissue)"),
        ("5   Label-swap permutation envelope", TEAL, "n_sim = 99 / 199  +  edge correction"),
        ("6   Effect-size (SES) ranking", PURPLE, "binary test  or  continuous SES"),
        ("OUTPUTS — networks · UMAP · heatmaps", BLUE_OUT, None),
    ]
    cur = y
    for i, (label, fill, chip) in enumerate(stages):
        box = rect(x, cur, bw, bh, fill, rounded=True)
        _shape_text(box, label, size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        if chip:
            ch = rect(chip_x, cur + (bh - chip_h) / 2, chip_w, chip_h, OBG,
                      rounded=True, line=OED)
            _shape_text(ch, chip, size=11, bold=False, color=OTX, italic=True)
        if i < len(stages) - 1:
            nxt = stages[i + 1][1]
            acol = TEAL if (fill == TEAL and nxt == TEAL) else PURPLE
            ar = shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Cm(cxc - 0.35),
                                  Cm(cur + bh + 0.02), Cm(0.7), Cm(gap - 0.05))
            ar.fill.solid(); ar.fill.fore_color.rgb = acol
            _no_line(ar); ar.shadow.inherit = False
        cur += bh + gap

    # quick-start panel
    cur += 0.15
    qh = 3.8
    rect(x, cur, w, qh, QBG, rounded=True, line=META)
    textbox(x + 0.5, cur + 0.22, w - 1.0, 1.0,
            [{"text": "Quick start — use NoSegs when…", "size": 15,
              "bold": True, "color": HEADING}])
    uses = [
        "segmentation fails or is unavailable on your sample",
        "you want ligand–receptor co-localisation straight from transcripts",
        "imaging platforms (CosMx, MERFISH, Xenium); ≥ ~50 transcripts / gene / ROI",
    ]
    ty = cur + 1.25
    for u in uses:
        textbox(x + 0.55, ty, 0.7, 0.55,
                [{"runs": [("✓", {"font": "DejaVu Sans", "size": 12.5,
                                    "bold": True, "color": GRN})]}])
        textbox(x + 1.2, ty, w - 1.8, 0.7,
                [{"text": u, "size": 12.5, "color": BODY}])
        ty += 0.8
    return cur + qh + 0.25


def heading(text, x, w, y, size=SZ_HEAD):
    textbox(x, y, w, 1.6, [{"text": text, "size": size, "bold": True,
                             "color": HEADING}])
    rect(x, y + 1.55, w, 0.07, PURPLE)        # thin rule under heading
    return y + 1.9


def subhead(text, x, w, y, size=SZ_SUBHEAD):
    textbox(x, y, w, 1.2, [{"text": text, "size": size, "bold": True,
                             "color": SUBHEAD}])
    return y + 1.25


def legend(tag, captext, x, w, y, size=SZ_LEGEND, gap=0.3):
    paras = [{
        "runs": [(tag + " ", {"bold": True, "size": size, "color": HEADING}),
                 (captext, {"size": size, "color": LEGEND})],
        "line_spacing": 1.02,
    }]
    h = est_block([{"text": tag + " " + captext, "size": size,
                    "line_spacing": 1.02}], w, size) + 0.15
    textbox(x, y, w, h + 0.4, paras)
    return y + h + gap


# ================= HEADER =================
rect(0, 0, 84.1, 12.4, PURPLE)
textbox(2.0, 1.0, 70, 5.2, [{
    "text": "NoSegs: Segmentation-Free Ligand–Receptor Co-Localisation in Spatial Transcriptomics",
    "size": 58, "bold": True, "color": WHITE, "line_spacing": 1.0}])
textbox(2.0, 6.5, 72, 1.7, [{
    "text": "Bivariate point-pattern statistics on raw transcripts — no cells, no clusters, no segmentation.",
    "size": 26, "italic": True, "color": META}])
textbox(2.0, 8.5, 76, 1.3, [{
    "text": "Harry Woodward   ·   Supervisor: Dr Dan Nicolau, Nicolau Lab   ·   MSc Bioinformatics, King's College London   ·   2026",
    "size": 22, "color": WHITE}])
textbox(2.0, 10.0, 76, 1.5, [{
    "text": "Platform: NanoString CosMx SMI    ·    Tissue: severe asthmatic, RSV-infected human lung    ·    ~2.66 M transcripts    ·    1,000-gene panel",
    "size": 18, "color": META}])
# KCL logo top-right
_lw = 8.6
shapes.add_picture(str(ASSETS / "pic_48.png"), Cm(84.1 - _lw - 1.6),
                   Cm(1.1), width=Cm(_lw))

# ================= COLUMN PANELS =================
COL_TOP, COL_BOT = 13.2, 100.6
col1, col2, col3 = 2.0, 29.2, 56.4
colw = 25.7
for cx in (col1, col2, col3):
    rect(cx, COL_TOP, colw, COL_BOT - COL_TOP, PANEL, rounded=True)


def report(colname, y):
    used = y - COL_TOP
    avail = COL_BOT - COL_TOP
    flag = "  <<< OVERFLOW" if y > COL_BOT else ""
    print(f"  {colname:7s} final y = {y:6.2f} cm  (used {used:5.2f} / {avail:.1f}){flag}")


# ================= LEFT COLUMN =================
x, w = col1 + 0.6, colw - 1.2
figw = colw - 1.2
y = COL_TOP + 0.5

y = heading("Introduction", x, w, y)
y = flow_text(x, w, y, [
    {"text": "Spatial transcriptomics maps gene expression at single-molecule resolution in intact tissue (He 2022). Almost every downstream analysis — cell-type calling, ligand–receptor inference, niche detection — requires a fragile first step: assigning each transcript to a cell.",
     "size": SZ_BODY, "space_after": 7, "line_spacing": 1.04},
    {"text": "On the CosMx sample analysed here, vendor segmentation leaves 43.6% of transcripts unassigned and most subcellular labels missing — segmentation quality is a known, platform-wide weakness (Petukhov 2022; Marco Salas 2025). Third-party tools (Cellpose, MOSAIK) cannot recover the signal from poor DAPI staining, so segmentation-dependent tools (CellChat, LIANA, Bento) become unusable (Jin 2021).",
     "size": SZ_BODY, "space_after": 7, "line_spacing": 1.04},
    {"text": "NoSegs takes a different route: treat each gene as a 2-D point pattern in tissue space and use bivariate Ripley's K-functions with label-permutation envelopes to detect ligand–receptor co-localisation directly. No cells. No clusters. Just transcripts.",
     "size": SZ_BODY, "space_after": 7, "line_spacing": 1.04},
    {"text": "The tissue is severe asthmatic lung during RSV infection — a clinically important setting where the spatial organisation of antiviral and immune signalling matters most, yet exactly where segmentation is least reliable.",
     "size": SZ_BODY, "line_spacing": 1.04},
], SZ_BODY, gap=0.5)

y = heading("Segmentation Loses Half the Signal", x, w, y)
y = flow_text(x, w, y, [
    {"text": "Vendor segmentation is unreliable on this sample. In a single field of view the cell mask captures only a fraction of the transcripts (Fig 1), and the loss is systematic across the dataset — 43.6% of all transcripts are never assigned to any cell (Fig 2). Every cell-level tool silently discards them.",
     "size": SZ_BODY, "line_spacing": 1.04}], SZ_BODY)
y = figure(FIGS / "14_M0_raw_data.png", col1 + 0.6, figw, y)
y = legend("Fig 1.",
           "One field of view, three layers. Left: raw DAPI morphology. Middle: the vendor cell-segmentation mask (148 cells) — the layer NoSegs discards. Right: the raw transcripts, with ~40% never assigned to any cell (magenta). NoSegs carries every point through.",
           x, w, y)
y = figure(FIGS / "14_M1_segmentation_dropout.png", col1 + 0.6, figw, y)
y = legend("Fig 2.",
           "Across all FOVs, vendor segmentation never assigns 43.6% of transcripts to a cell, so every cell-level tool silently discards them — consistently across the whole dataset.",
           x, w, y, gap=0.5)

y = heading("Preparing the Point Field", x, w, y)
y = flow_text(x, w, y, [
    {"text": "NoSegs keeps every transcript and works on the raw point cloud. Two light-touch steps prepare it: a per-FOV Gaussian mixture separates the infected and control tissue strips (Fig 3), and an adaptively-tuned DBSCAN pass removes diffuse background while retaining 87% of points (Fig 4).",
     "size": SZ_BODY, "line_spacing": 1.04}], SZ_BODY)
y = figure(FIGS / "14_M2_gmm_strips.png", col1 + 0.6, figw, y)
y = legend("Fig 3.",
           "After PCA aligns the tissue axis, a per-FOV 3-component Gaussian mixture separates the three tissue strips; boundaries are set per FOV, with manual overrides where the mixture is ambiguous.",
           x, w, y, gap=0.45)
y = subhead("Adaptive noise removal", x, w, y)
y = figure(FIGS / "14_M3_dbscan_qc.png", col1 + 0.6, figw, y)
y = legend("Fig 4.",
           "DBSCAN with an adaptively-set ε (97th-percentile 1-NN, clipped 20–30 px) flags diffuse background and stray clusters; 87% of transcripts survive QC.",
           x, w, y)
report("LEFT", y)

# ================= MIDDLE COLUMN =================
x, w = col2 + 0.6, colw - 1.2
figw_mid = 18.5          # compressed width for the statistics figures
y = COL_TOP + 0.5

y = heading("The NoSegs Pipeline, End to End", x, w, y)
y = flow_text(x, w, y, [
    {"text": "Eight stages take raw transcripts to signalling networks with no cell calling. Each gene becomes a 2-D point pattern; bivariate Ripley's K with a label-swap null then tests every ligand–receptor pair for spatial co-localisation. Teal marks the statistical core; amber chips mark the tunable choice at each step.",
     "size": SZ_BODY_M, "line_spacing": 1.03}], SZ_BODY_M, gap=0.3)
y = flowchart(x, y, w)
y = legend("Fig 5.",
           "The end-to-end workflow and quick-start guide. Eight native, editable stages from raw transcripts to signalling networks.",
           x, w, y, size=SZ_LEGEND_M, gap=0.45)

y = subhead("Observation window", x, w, y)
y = figure(FIGS / "14_M4_window_iterations.png", col2 + 0.6, colw - 1.2, y,
           target_w=figw_mid)
y = legend("Fig 6.",
           "The observation window is the dominant lever for inferential power. A concave hull tracks the true tissue extent, shrinking the area ~20× vs the bounding box and tightening the permutation null.",
           x, w, y, size=SZ_LEGEND_M, gap=0.4)

y = subhead("Multi-scale r-grid & edge correction", x, w, y)
y = figure(FIGS / "14_M5_rscale.png", col2 + 0.6, colw - 1.2, y,
           target_w=figw_mid)
y = legend("Fig 7.",
           "A deliberate multi-scale r-grid resolves both cell-scale and tissue-scale co-association; Shapely edge correction down-weights discs that spill outside the tissue window, following standard second-order point-pattern theory (Ripley 1976; Diggle 2003; Baddeley 2015).",
           x, w, y, size=SZ_LEGEND_M, gap=0.5)

y = heading("Validation & Controls", x, w, y)
y = flow_text(x, w, y, [
    {"text": "The method must first pass known-answer controls. A focal positive pair scores high and a diffuse housekeeping pair does not; both sit inside the permutation envelope under the production window, but their effect sizes separate cleanly (Figs 8–9). We therefore rank pairs by continuous standardised effect size (SES) — the continuous form of the global rank envelope test (Myllymäki 2017) — not binary significance.",
     "size": SZ_BODY_M, "line_spacing": 1.03}], SZ_BODY_M, gap=0.3)
y = figure(FIGS / "13_pos_vs_neg_sanity.png", col2 + 0.6, colw - 1.2, y,
           target_w=figw_mid)
y = legend("Fig 8.",
           "Both the positive control (KRT8×KRT18) and the housekeeping×specific negative (MALAT1×KRT18) fall inside the label-swap envelope on the infected strip, yet their effect sizes separate by 2.67 envelope half-widths. We therefore rank by continuous peak SES, not binary pass/fail.",
           x, w, y, size=SZ_LEGEND_M, gap=0.4)

# validation table (native editable table)
tbl_rows = [
    ("Control", "Pair", "Expected", "Observed"),
    ("Positive", "KRT8 × KRT18", "Above null", "SES = −2.18 (top)"),
    ("Negative", "MALAT1 × KRT18", "Far below", "SES = −4.85"),
]
tb_h = 4.0
gtbl = shapes.add_table(len(tbl_rows), 4, Cm(x), Cm(y), Cm(w), Cm(tb_h)).table
gtbl.columns[0].width = Cm(5.2)
gtbl.columns[1].width = Cm(7.6)
gtbl.columns[2].width = Cm(5.4)
gtbl.columns[3].width = Cm(w - 18.2)
for ri, row in enumerate(tbl_rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.margin_left = Cm(0.2); cell.margin_right = Cm(0.1)
        cell.margin_top = Cm(0.05); cell.margin_bottom = Cm(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.name = F
        r.font.size = Pt(14.5)
        r.font.bold = (ri == 0)
        if ri == 0:
            r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        else:
            r.font.color.rgb = BODY
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else PANEL
y += tb_h + 0.45

y = figure(FIGS / "13_C_group_ses_distribution.png", col2 + 0.6, colw - 1.2, y,
           target_w=figw_mid)
y = legend("Fig 9.",
           "Method is correctly null for diffuse markers (cytokeratins, housekeeping anchors) and elevated for focal pairs (paracrine, infection, stromal, immune). Bivariate K under this null detects excess spatial co-association beyond marginal density.",
           x, w, y, size=SZ_LEGEND_M)
report("MIDDLE", y)

# ================= RIGHT COLUMN =================
x, w = col3 + 0.6, colw - 1.2
figw = colw - 1.2
y = COL_TOP + 0.5

y = heading("Results — Infection Biology", x, w, y)
# bold take-home banner + big-number callout (editable)
band_h = 3.0
rect(x, y, w, band_h, PURPLE, rounded=True)
textbox(x + 4.8, y + 0.2, w - 5.2, band_h - 0.4, [
    {"text": "Immune cell-communication pairs preferentially co-localise on "
             "infected tissue — recovered with no segmentation step.",
     "size": 16.5, "bold": True, "color": WHITE, "line_spacing": 1.03,
     "space_after": 2},
    {"runs": [("Controls separate by ", {"size": 13.5, "color": META}),
              ("+2.67 envelope half-widths", {"size": 13.5, "bold": True, "color": WHITE}),
              (".", {"size": 13.5, "color": META})]},
])
textbox(x + 0.2, y + 0.15, 4.4, band_h - 0.3, [
    {"text": "5/6", "size": 42, "bold": True, "color": WHITE,
     "align": PP_ALIGN.CENTER, "space_after": 0, "line_spacing": 0.9},
    {"text": "immune pairs\npositive on the\ninfected strip", "size": 10,
     "color": META, "align": PP_ALIGN.CENTER, "line_spacing": 0.95},
])
y += band_h + 0.4
y = flow_text(x, w, y, [
    {"text": "Antigen presentation, T-cell recruitment and antiviral monocyte signalling all preferentially co-localise on the infected strip — the textbook adaptive immune response to RSV (Openshaw 2017; Russell 2017), recovered here with no segmentation step.",
     "size": SZ_BODY, "line_spacing": 1.04}], SZ_BODY, gap=0.45)

y = subhead("Infection signalling boost", x, w, y)
y = figure(FIGS / "13_F_infection_signature.png", col3 + 0.6, figw, y,
           target_w=20.0)
y = legend("Fig 10.",
           "Per-pair Strip-2 boost for the surviving infection-response pairs. Immune cell-communication pairs are positive (5/6): antigen presentation onto T cells (HLA-DRA×CD3D, +0.74; Roche & Furuta 2015), chemokine-driven T-cell recruitment (CXCL10×CD3D, +0.59; Groom & Luster 2011) and macrophage antigen presentation (HLA-DRA×CD68, +0.32) — the expected adaptive response to RSV (Openshaw 2017). A paradoxical second finding only a segmentation-free method can surface: antiviral interferon-stimulated genes paired with epithelium (ISG15×KRT8, MX1×KRT5; Schoggins 2011) move off infected tissue, consistent with epithelial loss redistributing antiviral activity into immune cells.",
           x, w, y, gap=0.45)

y = subhead("Co-localisation networks", x, w, y)
y = figure(FIGS / "13_G_network.png", col3 + 0.6, figw, y)
y = legend("Fig 11.",
           "Top-50 spatial co-association edges per strip with Louvain community detection. Communities are spatial signalling modules — interferon response, chemokine signalling, stromal ECM — recovered without any cell segmentation. Strip 2 (infected) carries more and larger modules than the controls.",
           x, w, y, gap=0.45)

y = subhead("Spatial signatures (UMAP)", x, w, y)
y = figure(FIGS / "13_H_umap.png", col3 + 0.6, figw, y, target_w=17.0)
y = legend("Fig 12.",
           "Each (pair, strip) treated as a 50-dim L(r) curve-shape vector, embedded with UMAP. Stromal, immune and infection-response pairs occupy distinct regions — the method recovers biological grouping as soft visual structure from raw transcripts alone.",
           x, w, y, gap=0.5)

y = heading("Where NoSegs Sits", x, w, y)
y = flow_text(x, w, y, [
    {"text": "No existing tool tests ligand–receptor co-localisation on raw transcripts. Cell-level methods require segmentation or cell-typing first; the one segmentation-free neighbour, FICTURE, infers tissue domains rather than pairwise interactions. NoSegs is the only method that satisfies all four capabilities at once (Fig 13).",
     "size": SZ_BODY, "line_spacing": 1.04}], SZ_BODY, gap=0.35)
y = figure(FIGS / "14_T_tool_positioning.png", col3 + 0.6, figw, y,
           target_w=19.0)
y = legend("Fig 13.",
           "No published tool screens L–R co-localisation on raw transcripts without segmentation. CellChat/LIANA (Jin 2021), Squidpy (Palla 2022), Bento (Mah 2024) and Baysor (Petukhov 2022) require segmentation or cell-typing; FICTURE (Si 2024) is segmentation-free but answers a different question (tissue domains); spatstat is a generic point-pattern library. NoSegs satisfies all four capabilities.",
           x, w, y)
report("RIGHT", y)

# ================= BOTTOM BAND =================
BB_TOP = 101.0
rect(0, BB_TOP, 84.1, 15.6, PURPLE)
rect(0, BB_TOP, 84.1, 0.12, META)
# Conclusions + next steps
textbox(2.0, BB_TOP + 0.5, 31, 0.9,
        [{"text": "Conclusions & Next Steps", "size": 22, "bold": True, "color": WHITE}])
textbox(2.0, BB_TOP + 1.7, 31, 9.0, [
    {"text": "Bivariate Ripley's K on raw transcript point clouds detects L–R co-localisation with no cell segmentation step.",
     "size": 15, "color": WHITE, "bullet": True, "space_after": 5, "line_spacing": 1.03},
    {"text": "Discriminates positive from negative controls by effect size (separation = +2.67 envelope half-widths under the production window).",
     "size": 15, "color": WHITE, "bullet": True, "space_after": 5, "line_spacing": 1.03},
    {"text": "Fills a vacant niche: Bento, Squidpy, CellChat, LIANA all require segmentation — NoSegs does not.",
     "size": 15, "color": WHITE, "bullet": True, "space_after": 10, "line_spacing": 1.03},
    {"runs": [("Next: ", {"bold": True, "size": 15, "color": META}),
              ("extended cytokeratin/housekeeping anchor panel; marker-gene pseudo-segmentation for cell-type stratification; benchmark vs Squidpy gr.ripley on cell-level reductions.",
               {"size": 15, "color": WHITE})],
     "line_spacing": 1.03},
])
# Key references (all 16, single column)
textbox(35.5, BB_TOP + 0.5, 30, 0.9,
        [{"text": "Key References", "size": 22, "bold": True, "color": WHITE}])
refs = [
    "Ripley BD (1976) J. Appl. Prob. 13:255–266 — the K-function.",
    "Diggle PJ (2003) Statistical Analysis of Spatial Point Patterns, 2nd ed. Arnold.",
    "Baddeley A, Rubak E, Turner R (2015) Spatial Point Patterns. CRC Press.",
    "Myllymäki M et al. (2017) J. R. Stat. Soc. B 79:381–404 — rank envelope test.",
    "He S et al. (2022) Nat. Biotechnol. 40:1794–1806 — CosMx SMI platform.",
    "Petukhov V et al. (2022) Nat. Biotechnol. 40:345–354 — Baysor segmentation.",
    "Si Y et al. (2024) Nat. Methods 21:1843–1854 — FICTURE (segmentation-free).",
    "Marco Salas S et al. (2025) Nat. Methods 22:813–823 — Xenium QA/segmentation.",
    "Palla G et al. (2022) Nat. Methods 19:171–178 — Squidpy.",
    "Mah CK et al. (2024) Genome Biol. 25:82 — Bento.",
    "Jin S et al. (2021) Nat. Commun. 12:1088 — CellChat / CellChatDB.",
    "Openshaw PJM et al. (2017) Annu. Rev. Immunol. 35:501–532 — RSV immunity.",
    "Russell CD et al. (2017) Clin. Microbiol. Rev. 30:481–502 — human RSV response.",
    "Schoggins JW et al. (2011) Nature 472:481–485 — ISG antiviral effectors.",
    "Roche PA, Furuta K (2015) Nat. Rev. Immunol. 15:203–216 — MHC class II.",
    "Groom JR, Luster AD (2011) Immunol. Cell Biol. 89:207–215 — CXCR3 ligands.",
]
textbox(35.5, BB_TOP + 1.7, 30, 13.0,
        [{"runs": [(f"{i+1}. ", {"bold": True, "size": 11.5, "color": META}),
                   (r, {"size": 11.5, "color": WHITE})],
          "space_after": 2, "line_spacing": 1.0}
         for i, r in enumerate(refs)])
# Acknowledgements
textbox(67.0, BB_TOP + 0.5, 13.5, 0.9,
        [{"text": "Acknowledgements", "size": 22, "bold": True, "color": WHITE}])
textbox(67.0, BB_TOP + 1.7, 13.0, 7.5, [
    {"text": "Dataset and HPC compute courtesy of King's College London (CREATE HPC, msc_appbio partition). We thank the Nicolau Lab and clinical collaborators for access to the RSV-infected lung tissue. CellChatDB via Jin et al. 2021. Project supervised by Dr Dan Nicolau, Nicolau Lab.",
     "size": 13, "color": WHITE, "line_spacing": 1.05}])
# QR code bottom-right
shapes.add_picture(str(ASSETS / "pic_50.png"), Cm(76.6), Cm(BB_TOP + 9.6),
                   width=Cm(5.4))

# ================= FOOTER BAR =================
rect(0, 116.9, 84.1, 2.0, RGBColor(0x3A, 0x05, 0x57))
textbox(2.0, 117.15, 80, 1.4,
        [{"text": "github.com/harrydaw/cosmx-pointpattern    ·    harrywoodward99@gmail.com    ·    MSc Bioinformatics — King's College London — 2026",
          "size": 14, "color": META, "align": PP_ALIGN.CENTER}])

prs.save(str(OUT))
print("saved", OUT)
print("slide cm:", round(prs.slide_width / 360000, 2),
      round(prs.slide_height / 360000, 2))
